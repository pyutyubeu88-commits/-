from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color System ────────────────────────────────────────────────────────────────
DARK   = RGBColor(0x07, 0x0D, 0x1B)   # near-black (impact slides)
DARK_C = RGBColor(0x0D, 0x18, 0x2C)   # card on dark bg
DARK_S = RGBColor(0x13, 0x22, 0x3E)   # subtle stripe on dark
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD   = RGBColor(0xF3, 0xF7, 0xFD)   # card bg (no border, contrast only)
STRIPE = RGBColor(0xE9, 0xF0, 0xF8)   # table row stripe
RED    = RGBColor(0xC8, 0x10, 0x2E)   # single accent color
PTXT   = RGBColor(0x08, 0x12, 0x24)   # primary text
STXT   = RGBColor(0x52, 0x66, 0x80)   # secondary text
TTXT   = RGBColor(0x8C, 0xA2, 0xBB)   # tertiary / labels
LINE   = RGBColor(0xD8, 0xE4, 0xEF)   # divider
GREEN  = RGBColor(0x08, 0x8C, 0x5C)   # success
RED_BG = RGBColor(0xFD, 0xEB, 0xEE)
GRN_BG = RGBColor(0xEB, 0xFB, 0xF4)

W, H  = Inches(13.33), Inches(7.5)
M     = Inches(0.5)
IG    = Inches(0.18)
CY    = Inches(1.10)    # content top
CB    = Inches(7.30)    # content bottom
CH    = CB - CY         # 6.20"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

# ── Primitives ──────────────────────────────────────────────────────────────────
def sl(): return prs.slides.add_slide(BLANK)

def R(s, l, t, w, h, fill=None, line=None, lw=Pt(0.5)):
    sh = s.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = lw
    else: sh.line.fill.background()
    return sh

def OV(s, l, t, w, h, c):
    sh = s.shapes.add_shape(9, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.fill.background(); return sh

def T(s, text, l, t, w, h, sz=11, bold=False, color=PTXT,
      align=PP_ALIGN.LEFT, italic=False):
    b = s.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return b

def TR(s, text, l, t, w, h, sz=11, bold=False, color=PTXT,
       align=PP_ALIGN.LEFT, fill=None, line=None, lw=Pt(0.5)):
    sh = R(s, l, t, w, h, fill=fill, line=line, lw=lw)
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    return sh

# ── Layout components ────────────────────────────────────────────────────────────
def hdr(s, num, title, sub=""):
    """Clean header: 3px red rule + small label + bold title + thin line."""
    R(s, 0, 0, W, Inches(0.03), fill=RED)
    T(s, num, M, Inches(0.07), Inches(0.5), Inches(0.22),
      sz=9, bold=True, color=RED)
    T(s, title, M, Inches(0.26), W - M*2, Inches(0.58),
      sz=28, bold=True, color=PTXT)
    if sub:
        T(s, sub, M, Inches(0.82), W - M*2, Inches(0.24),
          sz=10, color=STXT, italic=True)
    R(s, M, Inches(1.02), W - M*2, Inches(0.012), fill=LINE)

def pgn(s, n):
    T(s, f"{n} / 15", W - M - Inches(0.5), H - Inches(0.26),
      Inches(0.6), Inches(0.2), sz=8, color=TTXT, align=PP_ALIGN.RIGHT)

def card_lt(s, l, t, w, h):
    """Light card: CARD fill + 4px red left accent."""
    R(s, l, t, w, h, fill=CARD)
    R(s, l, t, Inches(0.04), h, fill=RED)

def card_dk(s, l, t, w, h):
    """Dark card: DARK_C fill + 3px red top accent."""
    R(s, l, t, w, h, fill=DARK_C)
    R(s, l, t, w, Inches(0.03), fill=RED)

def lbl(s, text, l, t):
    T(s, text.upper(), l, t, Inches(5), Inches(0.2),
      sz=8, bold=True, color=TTXT)

def sep(s, l, t, w):
    R(s, l, t, w, Inches(0.011), fill=LINE)

def blt(s, text, l, t, w, sz=11, c=PTXT):
    R(s, l, t + Inches(0.095), Inches(0.055), Inches(0.055), fill=RED)
    T(s, text, l + Inches(0.13), t, w - Inches(0.13), Inches(0.38),
      sz=sz, color=c)

def callout(s, l, t, w, h, body, tag=None):
    R(s, l, t, w, h, fill=DARK)
    R(s, l, t, Inches(0.04), h, fill=RED)
    if tag:
        T(s, tag, l + Inches(0.14), t + Inches(0.08),
          Inches(2), Inches(0.2), sz=8, bold=True, color=RED)
    T(s, body, l + Inches(0.14), t + (Inches(0.28) if tag else Inches(0.12)),
      w - Inches(0.22), h - (Inches(0.32) if tag else Inches(0.18)),
      sz=11.5, bold=True, color=WHITE)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S01  COVER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=DARK)
R(s, 0, 0, Inches(0.08), H, fill=RED)                          # left stripe
R(s, 0, H - Inches(0.05), W, Inches(0.05), fill=RED)           # bottom stripe

# decorative geometry
OV(s, W - Inches(6.5), Inches(-2.5), Inches(8.5), Inches(8.5), DARK_C)
OV(s, W - Inches(4.5), Inches(-0.5), Inches(4), Inches(4), DARK_S)
OV(s, Inches(1.0), H - Inches(3.8), Inches(4), Inches(4), DARK_C)

# eyebrow tag
TR(s, "소상공인 컨설팅 보고서",
   M + Inches(0.1), Inches(1.4), Inches(3.1), Inches(0.35),
   sz=11, bold=True, color=WHITE, fill=RED, align=PP_ALIGN.CENTER)

# hero title
T(s, "역삼효태권도장", M + Inches(0.1), Inches(1.92),
  Inches(11), Inches(1.5), sz=66, bold=True, color=WHITE)

# rule + subtitle
R(s, M + Inches(0.1), Inches(3.5), Inches(0.55), Inches(0.035), fill=RED)
T(s, "온라인 마케팅 진단  ·  자립형 마케팅 시스템 구축",
  M + Inches(0.1), Inches(3.64), Inches(10), Inches(0.46),
  sz=16, color=RGBColor(0x70, 0x90, 0xC0))

# meta block
meta = [
    ("업체명", "역삼효태권도장"),
    ("대표",   "김형열 관장"),
    ("팀",     "B그룹 5조  권용준 · 손미현"),
    ("기간",   "2026.05.26 – 2026.05.29  (3일  /  총 12시간)"),
    ("분야",   "학원업  |  온라인 마케팅 진단 및 홍보 콘텐츠 솔루션"),
]
R(s, M + Inches(0.1), Inches(4.3), Inches(6.2), Inches(2.42), fill=DARK_C)
R(s, M + Inches(0.1), Inches(4.3), Inches(0.04), Inches(2.42), fill=RED)
my = Inches(4.48)
for k, v in meta:
    T(s, k, M + Inches(0.26), my, Inches(0.95), Inches(0.35),
      sz=9, bold=True, color=RGBColor(0x60, 0x80, 0xB0))
    T(s, v, M + Inches(1.28), my, Inches(4.8), Inches(0.35),
      sz=9, color=RGBColor(0xBB, 0xCE, 0xE8))
    my += Inches(0.39)

pgn(s, 1)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S02  TOC
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "00", "목차", "Contents")
pgn(s, 2)

toc = [
    ("01", "업체 개요",           "역삼효태권도장 기본 정보 및 관원 현황"),
    ("02", "현황 진단 · 핵심 강점","온라인 마케팅 현황 분석 및 4대 강점 도출"),
    ("03", "1회차 컨설팅",         "브랜드 포지셔닝 및 온라인 전략 수립"),
    ("04", "2회차 컨설팅",         "네이버 플레이스 최적화  ·  AI 젬스 시스템 구축"),
    ("05", "3회차 컨설팅",         "자립 실습 완료  ·  리뷰 성장 계획  ·  총정리"),
    ("06", "컨설팅 성과 비교",     "Before → After  전후 비교"),
    ("07", "자립형 마케팅 루틴",   "주 2회 10분 루틴 및 실행 로드맵"),
    ("08", "향후 제언",            "지속가능성 및 성장 전략 제언"),
]
CW2 = (W - M*2 - IG) / 2
RH2 = (CH - IG*3) / 4

for i, (num, title, desc) in enumerate(toc):
    col, row = i % 2, i // 2
    x = M + col*(CW2 + IG)
    y = CY + row*(RH2 + IG)
    card_lt(s, x, y, CW2, RH2)
    # large number watermark
    T(s, num, x + CW2 - Inches(1.3), y + Inches(0.12),
      Inches(1.2), Inches(0.7), sz=38, bold=True,
      color=RGBColor(0xE2, 0xEC, 0xF8), align=PP_ALIGN.RIGHT)
    # small number
    T(s, num, x + Inches(0.18), y + Inches(0.18),
      Inches(0.55), Inches(0.35), sz=15, bold=True, color=RED)
    # title
    T(s, title, x + Inches(0.72), y + Inches(0.2),
      CW2 - Inches(0.9), Inches(0.38), sz=14, bold=True, color=PTXT)
    # desc
    T(s, desc, x + Inches(0.72), y + Inches(0.6),
      CW2 - Inches(0.9), Inches(0.3), sz=10, color=STXT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S03  업체 개요
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "01", "업체 개요", "역삼효태권도장 기본 정보 및 운영 현황")
pgn(s, 3)

# LEFT: 3 large stat blocks (stacked) ──────────────────────────────────────────
LW3 = Inches(4.2)
SH3 = (CH - IG*2) / 3

stats3 = [
    ("65", "명", "등록 관원"),
    ("12", "년", "운영 경력"),
    ("6",  "타임","일일 픽업 (직접 운행)"),
]
sy3 = CY
for val, unit, desc in stats3:
    card_lt(s, M, sy3, LW3, SH3)
    T(s, val, M + Inches(0.18), sy3 + Inches(0.2),
      Inches(1.8), Inches(0.9), sz=52, bold=True, color=PTXT)
    T(s, unit, M + Inches(0.18) + Inches(1.8), sy3 + Inches(0.44),
      Inches(0.9), Inches(0.46), sz=22, bold=False, color=STXT)
    T(s, desc, M + Inches(0.18), sy3 + SH3 - Inches(0.42),
      LW3 - Inches(0.3), Inches(0.3), sz=10, color=TTXT)
    sep(s, M + Inches(0.18), sy3 + SH3 - Inches(0.5), LW3 - Inches(0.3))
    sy3 += SH3 + IG

# RIGHT: info table ────────────────────────────────────────────────────────────
RX3 = M + LW3 + IG
RW3 = W - M - RX3
info3 = [
    ("명칭",     "역삼효태권도장"),
    ("위치",     "서울 강남구 역삼로 14길 18 2층\n역삼역 3번 출구 858m  ·  역삼초 인근"),
    ("영업시간", "월~금  12:00 – 20:00"),
    ("전화",     "02-552-7582"),
    ("운영인력", "관장 김형열  +  사범  총 2명"),
    ("등록인원", "65명  (유치부 16 · 저학년 23 · 고학년 23 · 중등 3)"),
    ("수강료",   "주 3~4회  월 180,000원  /  주 5회  월 190,000원"),
    ("운영기간", "12년  (인근 3년  +  현 위치 9년)"),
]
RH_ROW = CH / len(info3)
iy3 = CY
for i, (k, v) in enumerate(info3):
    bg = STRIPE if i % 2 == 0 else WHITE
    R(s, RX3, iy3, RW3, RH_ROW, fill=bg)
    T(s, k, RX3 + Inches(0.18), iy3 + Inches(0.08),
      Inches(1.2), RH_ROW, sz=9, bold=True, color=TTXT)
    T(s, v, RX3 + Inches(1.44), iy3 + Inches(0.08),
      RW3 - Inches(1.55), RH_ROW, sz=11, color=PTXT)
    iy3 += RH_ROW

R(s, RX3, CY, RW3, CH, line=LINE, lw=Pt(0.5))           # outer border only

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S04  현황 진단
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "02", "현황 진단", "온라인 마케팅 인프라 및 강약점 분석")
pgn(s, 4)

CALL_H4 = Inches(0.72)
CHIP_H4 = Inches(0.6)
PROB_H4 = CH - CALL_H4 - CHIP_H4 - IG*2
CW4 = (W - M*2 - IG) / 2

probs = [
    ("마케팅 인프라 부족",
     ["네이버 플레이스 리뷰 단 13개  (기본 4 / 키워드 7 / 블로그 2)",
      "브랜드명 혼재 — 역삼효태권도장 / 효태권도 병용",
      "채널별 정보 불일치로 검색 신뢰도 분산",
      "온라인 신규 유입 구조 사실상 전무"]),
    ("디지털 활용 역량 저하",
     ["인스타그램 등 SNS 채널 업데이트 미흡",
      "홍보 콘텐츠 자체 제작 경험 부족",
      "학부모 디지털 소통 루틴 부재",
      "마케팅 실행 루틴 전무"]),
]
px4 = M
for title, items in probs:
    card_lt(s, px4, CY, CW4, PROB_H4)
    lbl(s, "문제 진단", px4 + Inches(0.18), CY + Inches(0.12))
    T(s, title, px4 + Inches(0.18), CY + Inches(0.36),
      CW4 - Inches(0.26), Inches(0.42), sz=15, bold=True, color=PTXT)
    sep(s, px4 + Inches(0.18), CY + Inches(0.84), CW4 - Inches(0.26))
    ih4 = (PROB_H4 - Inches(0.98)) / len(items)
    iy4 = CY + Inches(0.98)
    for item in items:
        blt(s, item, px4 + Inches(0.18), iy4, CW4 - Inches(0.26), sz=11)
        iy4 += ih4
    px4 += CW4 + IG

callout(s, M, CY + PROB_H4 + IG, W - M*2, CALL_H4,
        "오프라인 서비스 품질·학부모 만족도는 매우 높으나, "
        "온라인 검색 결과와 신규 유입 구조로 연결되지 않고 있는 상태",
        tag="핵심 진단")

strengths = ["12년 장기 운영 신뢰", "역삼초 6타임 무료 픽업",
             "실시간 에듀패밀리 출결", "정서적 유대 교육 철학", "고학년 40% 유지율"]
sw = (W - M*2 - IG*4) / 5
sx4 = M
sy4 = CB - CHIP_H4
for st in strengths:
    TR(s, st, sx4, sy4, sw, CHIP_H4,
       sz=10.5, bold=True, color=PTXT, fill=CARD, align=PP_ALIGN.CENTER)
    sx4 += sw + IG

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S05  4대 핵심 강점  (dark)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=DARK)
R(s, 0, 0, Inches(0.08), H, fill=RED)
OV(s, W - Inches(5.5), Inches(-2), Inches(7), Inches(7), DARK_C)
pgn(s, 5)

T(s, "02", M + Inches(0.1), Inches(0.1), Inches(0.5), Inches(0.24),
  sz=9, bold=True, color=RED)
T(s, "4대 핵심 강점", M + Inches(0.1), Inches(0.3), Inches(10), Inches(0.82),
  sz=38, bold=True, color=WHITE)
T(s, "컨설팅 현장 조사를 통해 도출된 역삼효태권도장의 독보적 경쟁력",
  M + Inches(0.1), Inches(1.1), Inches(10), Inches(0.3),
  sz=12, color=RGBColor(0x70, 0x90, 0xC0))
R(s, M + Inches(0.1), Inches(1.44), Inches(0.5), Inches(0.03), fill=RED)

GY5  = Inches(1.6)
GH5  = H - Inches(0.15) - GY5
GCW5 = (W - M*2 - IG*3) / 4

cards5 = [
    ("01", "완벽한 케어",
     ["역삼초 6타임 무료 픽업\n(관장 직접 운행, 외주 無)",
      "실시간 에듀패밀리\n출결 사진 전송",
      "맞벌이 학부모 완벽 안심 케어"]),
    ("02", "모두가 주인공",
     ["시범단 위주가 아닌\n전원 성장 수업 구조",
      "예절·체력·태권도·\n생활지도 4축 체계",
      "60분 체계적 수업 설계"]),
    ("03", "통합 교육",
     ["기초체력 / 학교체육\n/ 태권도 / 생활지도",
      "정서적 유대 기반\n교육 철학 운영",
      "개인 맞춤 성장 지도"]),
    ("04", "12년의 신뢰",
     ["졸업 후에도 찾아오는\n검증된 교육력",
      "고학년·중등 비율 40%\n(장기 신뢰 객관 증거)",
      "\"가랑비에 옷 젖듯이\"\n서두르지 않는 철학"]),
]
gx5 = M
for num, title, buls in cards5:
    card_dk(s, gx5, GY5, GCW5, GH5)
    # large bg number
    T(s, num, gx5 + GCW5 - Inches(1.1), GY5 + Inches(0.06),
      Inches(1.0), Inches(0.78), sz=44, bold=True,
      color=DARK_S, align=PP_ALIGN.RIGHT)
    T(s, num, gx5 + Inches(0.18), GY5 + Inches(0.18),
      Inches(0.55), Inches(0.38), sz=14, bold=True, color=RED)
    T(s, title, gx5 + Inches(0.18), GY5 + Inches(0.62),
      GCW5 - Inches(0.28), Inches(0.52), sz=16, bold=True, color=WHITE)
    R(s, gx5 + Inches(0.18), GY5 + Inches(1.22),
      GCW5 - Inches(0.28), Inches(0.012),
      fill=RGBColor(0x22, 0x3A, 0x5E))
    by5 = GY5 + Inches(1.36)
    bh5 = (GH5 - Inches(1.36)) / len(buls)
    for b in buls:
        R(s, gx5 + Inches(0.18), by5 + Inches(0.09),
          Inches(0.05), Inches(0.05), fill=RED)
        T(s, b, gx5 + Inches(0.32), by5, GCW5 - Inches(0.42),
          bh5, sz=10.5, color=RGBColor(0xAA, 0xC2, 0xE0))
        by5 += bh5
    gx5 += GCW5 + IG

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S06  1회차 컨설팅
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "03", "1회차 컨설팅", "2026.05.26  —  브랜드 포지셔닝 및 온라인 전략 수립")
pgn(s, 6)

POS_H6  = Inches(0.78)
LOWER_H6= CH - POS_H6 - IG

callout(s, M, CY, W - M*2, POS_H6,
        "메인 포지셔닝  :  \"아이의 마음까지 책임지는  12년 신뢰\"",
        tag="확정된 포지셔닝")

LY6 = CY + POS_H6 + IG
LW6 = Inches(5.8)
RW6 = W - M*2 - LW6 - IG
RX6 = M + LW6 + IG

card_lt(s, M, LY6, LW6, LOWER_H6)
lbl(s, "브랜드 전략 확정", M + Inches(0.18), LY6 + Inches(0.12))
pos6 = [
    ("핵심 타깃",   "역삼동 맞벌이 학부모  (안심 케어 + 정서적 성장 동시 요구)"),
    ("브랜드 통일", "모든 채널  「역삼효태권도장 / 02-552-7582」  완전 일원화"),
    ("주요 채널",   "네이버 플레이스 최적화  +  인스타그램 주 2회 10분 루틴"),
    ("핵심 메시지", "\"가랑비에 옷 젖듯이\"  서두르지 않는 12년 교육 철학"),
]
rh6 = (LOWER_H6 - Inches(0.38)) / len(pos6)
iy6 = LY6 + Inches(0.38)
for k, v in pos6:
    T(s, k, M + Inches(0.18), iy6 + Inches(0.06),
      Inches(1.28), rh6, sz=9, bold=True, color=TTXT)
    T(s, v, M + Inches(1.52), iy6 + Inches(0.06),
      LW6 - Inches(1.66), rh6, sz=11.5, color=PTXT)
    sep(s, M + Inches(0.18), iy6 + rh6 - Inches(0.04), LW6 - Inches(0.26))
    iy6 += rh6

card_lt(s, RX6, LY6, RW6, LOWER_H6)
lbl(s, "1회차 수행 프로세스", RX6 + Inches(0.18), LY6 + Inches(0.12))
steps6 = [
    ("01", "온라인 현황 진단",
     "네이버 플레이스 현황  ·  브랜드명 혼재 및 채널 불일치"),
    ("02", "교육 철학 · 핵심 강점 도출",
     "12년 운영 신뢰도  ·  정서적 유대 기반 교육 정신 파악"),
    ("03", "4대 핵심 강점 정의",
     "완벽한 케어  /  모두가 주인공  /  통합 교육  /  12년의 신뢰"),
    ("04", "브랜드 포지셔닝 확정",
     "메인 포지셔닝·타깃 설정  ·  채널 전략 로드맵 수립·공유"),
]
sh6 = (LOWER_H6 - Inches(0.38)) / len(steps6)
sy6 = LY6 + Inches(0.38)
for num, title, desc in steps6:
    R(s, RX6 + Inches(0.18), sy6,
      Inches(0.58), sh6 - Inches(0.1), fill=DARK)
    T(s, num, RX6 + Inches(0.18), sy6 + (sh6 - Inches(0.1))*0.3,
      Inches(0.58), Inches(0.38), sz=14, bold=True, color=RED,
      align=PP_ALIGN.CENTER)
    R(s, RX6 + Inches(0.78), sy6,
      RW6 - Inches(0.96), sh6 - Inches(0.1), fill=STRIPE)
    T(s, title, RX6 + Inches(0.96), sy6 + Inches(0.06),
      RW6 - Inches(1.12), Inches(0.32), sz=12, bold=True, color=PTXT)
    T(s, desc, RX6 + Inches(0.96), sy6 + Inches(0.38),
      RW6 - Inches(1.12), sh6 - Inches(0.52), sz=10, color=STXT)
    sy6 += sh6

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S07  네이버 플레이스 최적화
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "04", "네이버 플레이스 전면 최적화",
    "2회차 컨설팅 2026.05.28  —  검색 최적화 완전 정비")
pgn(s, 7)

TIP_H7  = Inches(0.52)
THDR_H7 = Inches(0.42)
TROW_H7 = (CH - THDR_H7 - TIP_H7 - IG*2) / 5
TC7 = [M, Inches(2.62), Inches(5.5), Inches(9.62)]
TW7 = [Inches(2.17), Inches(2.83), Inches(4.07), W - M - Inches(9.62) - Inches(0.1)]

for j, (lbl7, cw) in enumerate(zip(
        ["구분", "변경 전", "변경 후  (최적화 완료)", "효과"], TW7)):
    TR(s, lbl7, TC7[j], CY, cw, THDR_H7,
       sz=11, bold=True, color=WHITE,
       fill=RED if j == 2 else DARK, align=PP_ALIGN.CENTER)

rows7 = [
    ("공식 상호\n및 번호", "상호명·번호 미통일",
     "역삼효태권도장 / 02-552-7582\n완전 통일", "검색 정확도↑\n브랜드 신뢰도↑"),
    ("영업 정보\n상태", "운영시간 미입력",
     "평일 12:00~20:00 입력\n상담 유도 시간 적용", "상담 전환율↑"),
    ("브랜드\n소개글", "단순 텍스트 / 매력 미흡",
     "\"가랑비에 옷 젖듯이...\" 12년 철학\n무료픽업·출결앱 키워드 반영", "감성 신뢰 형성"),
    ("시각 자료\n구성", "도장 사진 등록 부족",
     "수련·시설·픽업 사진\n신규 업로드 완료", "첫인상 신뢰도↑"),
    ("핵심 서비스\n노출", "핵심 강점 기능 미노출",
     "무료픽업·야간수련·에듀패밀리\n안심케어·체력강화 등록", "검색 노출↑"),
]
ty7 = CY + THDR_H7
for i, (c0, c1, c2, c3) in enumerate(rows7):
    bg7 = WHITE if i % 2 == 0 else STRIPE
    for j, (txt, cw) in enumerate(zip([c0, c1, c2, c3], TW7)):
        fc7 = PTXT if j != 2 else GREEN
        R(s, TC7[j], ty7, cw, TROW_H7,
          fill=GRN_BG if j == 2 else bg7, line=LINE, lw=Pt(0.4))
        T(s, txt, TC7[j] + Inches(0.1), ty7 + Inches(0.08),
          cw - Inches(0.14), TROW_H7 - Inches(0.1),
          sz=10.5, bold=(j == 0), color=fc7)
    ty7 += TROW_H7

callout(s, M, ty7 + IG, W - M*2, TIP_H7,
        "무료체험 예약 + 네이버 톡톡 기능 추가 → 예비 관원 소통 채널 확보 · 전방위 마케팅 활용 가능",
        tag="추가 완료")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S08  AI 젬스  (dark)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=DARK)
R(s, 0, 0, Inches(0.08), H, fill=RED)
OV(s, W - Inches(5), Inches(-1.5), Inches(6.5), Inches(6.5), DARK_C)
pgn(s, 8)

T(s, "04", M + Inches(0.1), Inches(0.1), Inches(0.5), Inches(0.24),
  sz=9, bold=True, color=RED)
T(s, "AI 젬스 자동화 시스템", M + Inches(0.1), Inches(0.3), Inches(11), Inches(0.82),
  sz=36, bold=True, color=WHITE)
T(s, "제미나이 젬스(Gemini Gems) 3종 + AI 홍보이미지 → 완전 자립형 마케팅 시스템",
  M + Inches(0.1), Inches(1.1), Inches(11), Inches(0.3),
  sz=12, color=RGBColor(0x70, 0x90, 0xC0))
R(s, M + Inches(0.1), Inches(1.44), Inches(0.5), Inches(0.03), fill=RED)

GY8  = Inches(1.62)
GH8  = H - Inches(0.15) - GY8
GCW8 = (W - M*2 - IG*3) / 4

gems8 = [
    ("GEMS 1", "인스타\n포스팅 생성기",
     "사진만 넣으면\n인스타 텍스트와\n해시태그 자동 생성",
     "주 2회  5분 완성"),
    ("GEMS 2", "리뷰 요청\n카톡 생성기",
     "학부모에게 부담 없는\n개인 맞춤 리뷰 요청\n메시지 자동 작성",
     "30개 달성 핵심 도구"),
    ("GEMS 3", "리뷰 답글\n생성기",
     "네이버 리뷰 입력 시\n전문적인 감사 답글\n자동 생성",
     "응답률 100%"),
    ("AI 이미지", "홍보 이미지\n생성",
     "챗GPT · 제미나이 활용\n아이 초상권 보호\n일러스트 제작",
     "홍보 외주비 절감"),
]
gx8 = M
for badge, title, desc, result in gems8:
    card_dk(s, gx8, GY8, GCW8, GH8)
    # badge bar
    R(s, gx8, GY8, GCW8, Inches(0.42), fill=DARK_S)
    T(s, badge, gx8 + Inches(0.18), GY8 + Inches(0.1),
      GCW8 - Inches(0.24), Inches(0.26), sz=11, bold=True, color=RED)
    T(s, title, gx8 + Inches(0.18), GY8 + Inches(0.56),
      GCW8 - Inches(0.26), Inches(0.84), sz=16, bold=True, color=WHITE)
    R(s, gx8 + Inches(0.18), GY8 + Inches(1.52),
      GCW8 - Inches(0.28), Inches(0.012),
      fill=RGBColor(0x1E, 0x32, 0x56))
    T(s, desc, gx8 + Inches(0.18), GY8 + Inches(1.68),
      GCW8 - Inches(0.26), GH8 - Inches(2.58),
      sz=11, color=RGBColor(0xA0, 0xBE, 0xDC))
    # result footer
    R(s, gx8, GY8 + GH8 - Inches(0.88), GCW8, Inches(0.88), fill=DARK_S)
    R(s, gx8, GY8 + GH8 - Inches(0.88), GCW8, Inches(0.03), fill=RED)
    T(s, result, gx8 + Inches(0.18), GY8 + GH8 - Inches(0.78),
      GCW8 - Inches(0.26), Inches(0.65),
      sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    gx8 += GCW8 + IG

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S09  관원 구성 분석
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "04", "관원 구성 데이터 분석",
    "출석부 기반 연령별 분석 → 마케팅 타깃 전략 반영")
pgn(s, 9)

LW9 = Inches(7.0)
RW9 = W - M*2 - LW9 - IG
RX9 = M + LW9 + IG

# table ────────────────────────────────────────────────────────────────────────
tc9 = [M, M + Inches(1.65), M + Inches(2.6), M + Inches(3.42)]
tw9 = [Inches(1.6), Inches(0.9), Inches(0.78), LW9 - Inches(3.45)]

for j, (h_, cw) in enumerate(zip(["연령대","인원","비율","마케팅 포인트"], tw9)):
    T(s, h_, tc9[j], CY, cw, Inches(0.28),
      sz=9, bold=True, color=TTXT)
sep(s, M, CY + Inches(0.28), LW9)

mrows9 = [
    ("7세",        "10명","15.4%","초등 입학 준비 전문성 어필"),
    ("1·2학년",    "12명","18.5%","기초 체력 형성기 집중 케어"),
    ("3학년",      "11명","16.9%","기초 체력 단단한 허리층 타깃"),
    ("4·5·6학년",  "29명","44.6%","학교 체육·교우관계 연계 부각"),
    ("중학교",      "3명", "4.6%","장기 신뢰 극대화 유지 전략"),
    ("합계",       "65명","100%", "역삼동 중심가 안정적 운영 증명"),
]
RH9 = (CH - Inches(0.28) - Inches(0.015)) / len(mrows9)
ty9 = CY + Inches(0.3)
for i, (age, cnt, pct, tip) in enumerate(mrows9):
    last = i == len(mrows9) - 1
    bg9 = DARK if last else (STRIPE if i % 2 == 0 else WHITE)
    for j, (tv, cw) in enumerate(zip([age, cnt, pct, tip], tw9)):
        fc9 = (WHITE if last else
               RED if j == 0 else PTXT)
        R(s, tc9[j], ty9, cw, RH9, fill=bg9,
          line=LINE if not last else None, lw=Pt(0.35))
        T(s, tv, tc9[j] + Inches(0.06), ty9 + Inches(0.06),
          cw - Inches(0.08), RH9 - Inches(0.08),
          sz=11, bold=(j == 0 or last), color=fc9,
          align=PP_ALIGN.CENTER if j < 3 else PP_ALIGN.LEFT)
    ty9 += RH9

# bar chart + insight ──────────────────────────────────────────────────────────
INSIGHT_H9 = Inches(1.82)
BAR_H9     = CH - INSIGHT_H9 - IG
bars9 = [("7세",15.4),("1·2학년",18.5),("3학년",16.9),("4~6학년",44.6),("중학교",4.6)]
br9   = BAR_H9 / len(bars9)
MAXBW = RW9 - Inches(0.18)
lbl(s, "연령별 비율", RX9, CY + Inches(0.04))
by9 = CY + Inches(0.26)
for lbl9, pct9 in bars9:
    T(s, lbl9, RX9, by9 + br9*0.08, Inches(1.28), br9*0.6,
      sz=9.5, color=STXT)
    bw9 = (MAXBW - Inches(1.3)) * pct9 / 50
    R(s, RX9 + Inches(1.32), by9 + br9*0.15,
      MAXBW - Inches(1.32), br9*0.55, fill=STRIPE)
    R(s, RX9 + Inches(1.32), by9 + br9*0.15,
      bw9, br9*0.55, fill=DARK)
    T(s, f"{pct9}%",
      RX9 + Inches(1.32) + bw9 + Inches(0.06), by9 + br9*0.08,
      Inches(0.6), br9*0.6, sz=9, bold=True, color=PTXT)
    by9 += br9

R(s, RX9, CY + BAR_H9 + IG, RW9, INSIGHT_H9, fill=DARK)
R(s, RX9, CY + BAR_H9 + IG, Inches(0.04), INSIGHT_H9, fill=RED)
T(s, "고학년 40% 유지율",
  RX9 + Inches(0.16), CY + BAR_H9 + IG + Inches(0.14),
  RW9 - Inches(0.24), Inches(0.36), sz=13, bold=True, color=RED)
T(s, "일반적으로 학과 공부로 이탈하는 연령대가\n전체의 40%를 차지  →  \"한 번 믿으면 끝까지\"\n장기 신뢰 도장의 가장 강력한 마케팅 근거",
  RX9 + Inches(0.16), CY + BAR_H9 + IG + Inches(0.54),
  RW9 - Inches(0.24), INSIGHT_H9 - Inches(0.64),
  sz=10.5, color=RGBColor(0xA0, 0xBE, 0xDC))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S10  3회차 컨설팅
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "05", "3회차 컨설팅",
    "2026.05.29  —  자립 실습 완료  ·  리뷰 성장 계획  ·  시스템 총정리")
pgn(s, 10)

LW10 = Inches(5.9)
RW10 = W - M*2 - LW10 - IG
RX10 = M + LW10 + IG

card_lt(s, M, CY, LW10, CH)
lbl(s, "관장님 자립 실습 완료", M + Inches(0.18), CY + Inches(0.12))
T(s, "5분 루틴 완전 체득 확인",
  M + Inches(0.18), CY + Inches(0.36), LW10 - Inches(0.26), Inches(0.38),
  sz=14, bold=True, color=PTXT)
sep(s, M + Inches(0.18), CY + Inches(0.82), LW10 - Inches(0.26))

practice10 = [
    ("젬스 3번  리뷰 답글 생성기 활용",
     "네이버 플레이스 기존 리뷰 답글 실제 등록 완료"),
    ("젬스 1번  인스타 포스팅 생성기 활용",
     "인스타그램 2번째 포스팅 업로드 · 품질 조정 완료"),
    ("5분 이내 완성 루틴 완전 체득",
     "사진(1분) → 젬스 문구(2분) → 업로드(2분) = 5분"),
    ("네이버 플레이스 보완 완료",
     "무료체험 예약 + 네이버 톡톡 기능 추가"),
    ("AI 홍보 이미지 생성 실습",
     "챗GPT·제미나이 활용 초상권 보호 일러스트 제작"),
]
pih = (CH - Inches(0.88)) / len(practice10)
py10 = CY + Inches(0.88)
for title, desc in practice10:
    TR(s, "✓", M + Inches(0.18), py10 + (pih - Inches(0.36))/2,
       Inches(0.36), Inches(0.36),
       sz=13, bold=True, color=WHITE, fill=GREEN, align=PP_ALIGN.CENTER)
    T(s, title, M + Inches(0.66), py10 + Inches(0.04),
      LW10 - Inches(0.8), pih*0.48, sz=11.5, bold=True, color=PTXT)
    T(s, desc, M + Inches(0.66), py10 + pih*0.5,
      LW10 - Inches(0.8), pih*0.46, sz=10, color=STXT)
    sep(s, M + Inches(0.18), py10 + pih - Inches(0.04), LW10 - Inches(0.26))
    py10 += pih

card_lt(s, RX10, CY, RW10, CH)
lbl(s, "리뷰 30개 달성 실행 계획", RX10 + Inches(0.18), CY + Inches(0.12))
T(s, "현재 13개 → 목표 30개  (3개월)",
  RX10 + Inches(0.18), CY + Inches(0.36), RW10 - Inches(0.26), Inches(0.36),
  sz=14, bold=True, color=PTXT)

# progress bar
PBY = CY + Inches(0.86)
R(s, RX10 + Inches(0.18), PBY, RW10 - Inches(0.26), Inches(0.36), fill=STRIPE)
R(s, RX10 + Inches(0.18), PBY,
  (RW10 - Inches(0.26)) * 13/30, Inches(0.36), fill=DARK)
T(s, "13개", RX10 + Inches(0.26), PBY + Inches(0.08),
  Inches(0.8), Inches(0.22), sz=9, bold=True, color=WHITE)
sep(s, RX10 + Inches(0.18), PBY + Inches(0.44), RW10 - Inches(0.26))

timings10 = [
    ("신규 등록 1개월 후",    "적응 완료 → 만족도 최고점"),
    ("승급 심사 직후",        "성취감 최고조 → 긍정 리뷰 자연 유도"),
    ("야간수련·이벤트 직후",  "특별 경험 → 감동 공유 욕구"),
    ("학부모 감사 표현 시",   "자발적 만족 → 최적 요청 타이밍"),
]
T(s, "최적 리뷰 요청 타이밍",
  RX10 + Inches(0.18), PBY + Inches(0.54), RW10 - Inches(0.26), Inches(0.22),
  sz=8, bold=True, color=TTXT)
tih = (CB - PBY - Inches(0.84)) / len(timings10)
ty10 = PBY + Inches(0.82)
for timing, reason in timings10:
    R(s, RX10 + Inches(0.18), ty10, Inches(0.04), tih - Inches(0.06), fill=RED)
    T(s, timing, RX10 + Inches(0.32), ty10 + Inches(0.04),
      RW10 - Inches(0.46), tih*0.48, sz=11.5, bold=True, color=PTXT)
    T(s, reason, RX10 + Inches(0.32), ty10 + tih*0.5,
      RW10 - Inches(0.46), tih*0.46, sz=10, color=STXT)
    sep(s, RX10 + Inches(0.18), ty10 + tih - Inches(0.04), RW10 - Inches(0.26))
    ty10 += tih

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S11  Before / After
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "06", "컨설팅 성과 비교", "Before → After  전후 비교")
pgn(s, 11)

BAHDR = Inches(0.42)
BAROW = (CH - BAHDR) / 4
BC = [M, Inches(2.68), Inches(7.34)]
BW = [Inches(2.23), Inches(4.61), W - M - Inches(7.34) - Inches(0.1)]

for j, (lbl_, cw, bg_) in enumerate(zip(
        ["항목", "BEFORE  컨설팅 이전", "AFTER  컨설팅 이후"],
        BW, [DARK, RGBColor(0x80,0x08,0x1C), GREEN])):
    TR(s, lbl_, BC[j], CY, cw, BAHDR,
       sz=11, bold=True, color=WHITE, fill=bg_, align=PP_ALIGN.CENTER)

ba_rows = [
    ("네이버 플레이스\n운영",
     "운영 미흡\n리뷰 댓글 無 · 소식 無 · 정보 불일치",
     "전면 정비 완료\n메뉴·콘텐츠·댓글·소식 업데이트"),
    ("마케팅\n문구 작성",
     "인스타/블로그 문구\n자체 작성 불가",
     "AI 젬스 3종 + 매뉴얼 제공\n클릭 한 번으로 자동 생성"),
    ("홍보물\n제작",
     "홍보물 자체 제작\n완전 불가",
     "챗GPT·제미나이 활용\n초상권 보호 일러스트 자유 제작"),
    ("마케팅\n실행 루틴",
     "마케팅 루틴 전무",
     "주 2회 10분 루틴 확립\n촬영1분 + 젬스2분 + 업로드2분"),
]
cy11 = CY + BAHDR
for i, (label, before, after) in enumerate(ba_rows):
    R(s, BC[0], cy11, BW[0], BAROW, fill=CARD, line=LINE, lw=Pt(0.4))
    T(s, label, BC[0] + Inches(0.1), cy11 + (BAROW - Inches(0.42))/2,
      BW[0] - Inches(0.14), Inches(0.42),
      sz=12, bold=True, color=PTXT, align=PP_ALIGN.CENTER)
    R(s, BC[1], cy11, BW[1], BAROW, fill=RED_BG, line=LINE, lw=Pt(0.4))
    T(s, before, BC[1] + Inches(0.14), cy11 + Inches(0.14),
      BW[1] - Inches(0.2), BAROW - Inches(0.2), sz=11, color=PTXT)
    T(s, "→", Inches(6.82), cy11 + (BAROW - Inches(0.38))/2,
      Inches(0.46), Inches(0.38), sz=20, bold=True,
      color=TTXT, align=PP_ALIGN.CENTER)
    R(s, BC[2], cy11, BW[2], BAROW, fill=GRN_BG, line=LINE, lw=Pt(0.4))
    T(s, after, BC[2] + Inches(0.14), cy11 + Inches(0.14),
      BW[2] - Inches(0.2), BAROW - Inches(0.2),
      sz=11, bold=True, color=GREEN)
    cy11 += BAROW

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S12  자립형 마케팅 루틴
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "07", "자립형 마케팅 루틴",
    "주 2회 10분 루틴  —  컨설팅 이후에도 지속 가능한 성장")
pgn(s, 12)

FLOW_H12 = Inches(2.56)
CALL_H12 = Inches(0.54)
CLK_H12  = CH - FLOW_H12 - CALL_H12 - IG*2

steps12 = [
    ("01\n사진 촬영","1분","수련 뒷모습\n시설 모습"),
    ("02\n젬스 문구","2분","인스타 텍스트\n해시태그 자동"),
    ("03\n업로드",  "2분","인스타그램\n포스팅 완성"),
    ("04\n리뷰 요청","수시","젬스 2번 활용\n맞춤 카톡"),
    ("05\n리뷰 답글","즉시","젬스 3번 활용\n자동 생성"),
]
SCW12 = (W - M*2 - IG*4) / 5
sx12 = M
for i, (label, time, desc) in enumerate(steps12):
    card_lt(s, sx12, CY, SCW12, FLOW_H12)
    R(s, sx12 + Inches(0.04), CY, SCW12 - Inches(0.04), Inches(0.4), fill=DARK)
    T(s, time, sx12 + Inches(0.18), CY + Inches(0.08),
      SCW12 - Inches(0.26), Inches(0.26), sz=13, bold=True, color=RED)
    T(s, label, sx12 + Inches(0.18), CY + Inches(0.5),
      SCW12 - Inches(0.26), Inches(0.78), sz=13, bold=True, color=PTXT)
    sep(s, sx12 + Inches(0.18), CY + Inches(1.36), SCW12 - Inches(0.26))
    for j, line in enumerate(desc.split("\n")):
        T(s, line, sx12 + Inches(0.18), CY + Inches(1.5) + j*Inches(0.42),
          SCW12 - Inches(0.26), Inches(0.4), sz=10.5, color=STXT)
    if i < 4:
        T(s, "›", sx12 + SCW12 + Inches(0.02), CY + FLOW_H12*0.4,
          Inches(0.16), Inches(0.4), sz=20, bold=True,
          color=TTXT, align=PP_ALIGN.CENTER)
    sx12 += SCW12 + IG

callout(s, M, CY + FLOW_H12 + IG, W - M*2, CALL_H12,
        "총 5분 이내 완성  —  스마트폰 하나로 전문적인 마케팅 콘텐츠 제작 가능")

chk12 = [
    ("네이버 플레이스",  ["정보 최신화","리뷰 답글 등록","소식 게시"]),
    ("인스타그램",       ["주 2회 포스팅","해시태그 최적화","팔로워 소통"]),
    ("리뷰 관리",        ["개별 요청 발송","답글 즉시 등록","30개 목표 진행"]),
    ("환경개선",         ["계단·복도 청결","비포·애프터 기록","학부모 첫인상 관리"]),
]
CCW = (W - M*2 - IG*3) / 4
ccx = M
ccy = CY + FLOW_H12 + IG + CALL_H12 + IG
for ctitle, citems in chk12:
    card_lt(s, ccx, ccy, CCW, CLK_H12)
    T(s, ctitle, ccx + Inches(0.18), ccy + Inches(0.12),
      CCW - Inches(0.26), Inches(0.32), sz=12, bold=True, color=PTXT)
    sep(s, ccx + Inches(0.18), ccy + Inches(0.5), CCW - Inches(0.28))
    iih = (CLK_H12 - Inches(0.56)) / len(citems)
    iy12 = ccy + Inches(0.56)
    for ci in citems:
        R(s, ccx + Inches(0.18), iy12 + iih*0.35,
          Inches(0.06), Inches(0.06), fill=GREEN)
        T(s, ci, ccx + Inches(0.34), iy12 + Inches(0.04),
          CCW - Inches(0.46), iih*0.8, sz=10.5, color=PTXT)
        iy12 += iih
    ccx += CCW + IG

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S13  자립 시스템 총정리
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "05", "자립 시스템 총정리", "컨설팅 종료 시점 구축 완료 현황")
pgn(s, 13)

items13 = [
    ("네이버 플레이스 최적화 및 업데이트",
     "브랜드 통일·운영시간·소개글·사진·핵심 서비스 전면 정비 / 무료체험 예약 + 네이버 톡톡 추가 → 예비 관원 소통 채널 확보"),
    ("인스타그램 공식 채널 정상 가동",
     "공식 인스타그램 개설·프로필 구성·첫 포스팅 완료 / 관장님 직접 실습으로 지속 운영 역량 확보"),
    ("AI 젬스 4종 완전 체득 및 실습 완료",
     "인스타 포스팅 생성기 / 리뷰 요청 카톡 생성기 / 리뷰 답글 생성기 / AI 홍보이미지 생성 → 매뉴얼 제작·숙지 확인"),
    ("주 2회 10분 마케팅 루틴 확립",
     "사진(1분) + 젬스(2분) + 업로드(2분) = 5분 루틴 / 외부 대행 없이 스스로 운영 가능한 자립형 마케팅 시스템"),
    ("리뷰 30개 달성 단계별 실행 계획 수립",
     "현재 13개 → 3개월 내 30개 목표 / 최적 타이밍별 요청 전략 / 젬스 2번 개별 맞춤 발송으로 자연스러운 리뷰 확보"),
    ("환경개선 및 학부모 눈높이 소통 의지 제고",
     "계단·복도 청결 환경개선 / 서울시 소상공인 지원사업 연계 시설개선 최대 300만원 지원 안내"),
]
IH13 = (CH - IG*5) / 6
iy13 = CY
for i, (title, desc) in enumerate(items13):
    bg13 = STRIPE if i % 2 == 0 else WHITE
    R(s, M, iy13, W - M*2, IH13, fill=bg13)
    R(s, M, iy13, Inches(0.04), IH13, fill=RED)
    TR(s, "✓", M + Inches(0.1), iy13 + (IH13 - Inches(0.44))/2,
       Inches(0.44), Inches(0.44),
       sz=14, bold=True, color=WHITE, fill=GREEN, align=PP_ALIGN.CENTER)
    T(s, title, M + Inches(0.66), iy13 + Inches(0.08),
      Inches(5.2), Inches(0.3), sz=12.5, bold=True, color=PTXT)
    T(s, desc, M + Inches(0.66), iy13 + Inches(0.4),
      W - M*2 - Inches(0.78), IH13 - Inches(0.48),
      sz=10, color=STXT)
    iy13 += IH13 + IG

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S14  향후 제언
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=WHITE)
hdr(s, "08", "향후 제언", "지속 가능성 및 성장 전략 제언")
pgn(s, 14)

QH14  = Inches(0.9)
CH14  = (CH - QH14 - IG*3) / 2
CW14  = (W - M*2 - IG) / 2

recs14 = [
    ("01", "단톡방 유지 사후 관리",
     "컨설팅 팀과 단체 채팅방 유지로 추후 궁금증 즉시 해결\n지속적인 관리 및 피드백 제공 예정"),
    ("02", "WOM 입소문 마케팅",
     "컨설팅 수혜 업체가 주변 소상공인 대상으로 홍보\n→ 추천제도 운영으로 선순환 생태계 구축 제언"),
    ("03", "서울시 지원사업 연계",
     "서울특별시 소상공인 지원사업 활용 (seoulsbdc.or.kr)\n하반기 시설개선 최대 300만원 지원 활용 제언"),
    ("04", "콘텐츠 자산 지속 축적",
     "수련 사진·영상 꾸준한 기록 → 인스타·네이버 자산화\n\"가랑비에 옷 젖듯이\" 철학을 온라인에서도 실현"),
]
for i, (num, title, desc) in enumerate(recs14):
    cx14 = M if i % 2 == 0 else M + CW14 + IG
    cy14 = CY if i < 2 else CY + CH14 + IG
    card_lt(s, cx14, cy14, CW14, CH14)
    T(s, num, cx14 + Inches(0.18), cy14 + Inches(0.16),
      Inches(0.52), Inches(0.52), sz=22, bold=True, color=RED)
    T(s, title, cx14 + Inches(0.72), cy14 + Inches(0.22),
      CW14 - Inches(0.88), Inches(0.36), sz=13.5, bold=True, color=PTXT)
    sep(s, cx14 + Inches(0.18), cy14 + Inches(0.64), CW14 - Inches(0.26))
    for j, line in enumerate(desc.split("\n")):
        T(s, line, cx14 + Inches(0.18), cy14 + Inches(0.8) + j*Inches(0.38),
          CW14 - Inches(0.26), Inches(0.36), sz=11, color=PTXT)

# quote
QY14 = CB - QH14
R(s, M, QY14, W - M*2, QH14, fill=DARK)
R(s, M, QY14, Inches(0.04), QH14, fill=RED)
T(s, "관장님 소회", M + Inches(0.16), QY14 + Inches(0.1),
  Inches(1.5), Inches(0.2), sz=8, bold=True, color=RED)
T(s, "\"태권도 학원들이 대체로 홍보 마케팅과 거리가 멀었는데, "
     "3일 동안 완전히! 꼭 필요한 것을 너무나도 꽉꽉 채워주셨습니다.  감사합니다!\"",
  M + Inches(0.16), QY14 + Inches(0.3),
  W - M*2 - Inches(0.24), Inches(0.52),
  sz=13, bold=True, color=WHITE)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# S15  CLOSING
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = sl()
R(s, 0, 0, W, H, fill=DARK)
R(s, 0, 0, Inches(0.08), H, fill=RED)
R(s, 0, H - Inches(0.05), W, Inches(0.05), fill=RED)
OV(s, W - Inches(6.8), Inches(-2.5), Inches(9), Inches(9), DARK_C)
OV(s, W - Inches(2.8), H - Inches(4.2), Inches(5.5), Inches(5.5), DARK_S)

T(s, "감사합니다", M + Inches(0.1), Inches(1.5),
  Inches(10), Inches(1.5), sz=68, bold=True, color=WHITE)
R(s, M + Inches(0.1), Inches(3.1), Inches(0.72), Inches(0.035), fill=RED)
T(s, "역삼효태권도장  온라인 마케팅 컨설팅 보고서",
  M + Inches(0.1), Inches(3.25), Inches(10), Inches(0.44),
  sz=15, color=RGBColor(0x68, 0x88, 0xB8))

summary15 = [
    "네이버 플레이스 완전 최적화 및 브랜드 통일",
    "인스타그램 공식 채널 개설 및 첫 포스팅 완료",
    "AI 젬스 3종 자동화 시스템 완전 체득",
    "주 2회 10분 마케팅 루틴 자립 시스템 완성",
    "3개월 리뷰 30개 달성 단계별 실행 계획 수립",
]
sy15 = Inches(3.92)
for pt in summary15:
    R(s, M + Inches(0.1), sy15 + Inches(0.1),
      Inches(0.04), Inches(0.14), fill=RED)
    T(s, pt, M + Inches(0.26), sy15, Inches(7.5), Inches(0.36),
      sz=12, color=RGBColor(0xBB, 0xCE, 0xE8))
    sy15 += Inches(0.44)

# team info
R(s, M + Inches(0.1), H - Inches(1.2), Inches(7.2), Inches(0.68),
  fill=DARK_C)
R(s, M + Inches(0.1), H - Inches(1.2), Inches(0.04), Inches(0.68), fill=RED)
T(s, "컨설팅 팀", M + Inches(0.24), H - Inches(1.12),
  Inches(1.2), Inches(0.22), sz=8, bold=True, color=RED)
T(s, "B그룹 5조  권용준 · 손미현  |  2026.05.26 – 2026.05.29",
  M + Inches(0.24), H - Inches(0.9), Inches(6.8), Inches(0.32),
  sz=11, color=RGBColor(0xBB, 0xCE, 0xE8))

pgn(s, 15)

# ── Save ────────────────────────────────────────────────────────────────────────
OUT = "/home/user/-/역삼효태권도장_온라인마케팅컨설팅_보고서.pptx"
prs.save(OUT)
print("Saved:", OUT)
